import os
import datetime
from flask import Flask, request, jsonify, render_template
import requests

app = Flask(__name__)

# ─────────────────────────────────────────────────────────────────────────────
#  FREE AI MODEL — Groq Cloud (llama-3.3-70b-versatile)
#  ✅ 100% Free  ✅ No billing required  ✅ Very fast
#
#  STEP 1: Go to https://console.groq.com
#  STEP 2: Sign up free → API Keys → Create Key
#  STEP 3: Paste the key below (replace the placeholder)
# ─────────────────────────────────────────────────────────────────────────────
GROQ_API_KEY = "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"   # ← PASTE YOUR KEY HERE
GROQ_MODEL   = "llama-3.3-70b-versatile"
GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
# ─────────────────────────────────────────────────────────────────────────────

SYSTEM_PROMPT = (
    "You are Arisu, a professional and intelligent AI assistant. "
    "Be precise, helpful, and articulate. Provide accurate, structured responses. "
    "If you do not know something, say so honestly. "
    "When the user shares file contents, analyse and respond based on that content. "
    "You are part of Arisu AI, a product under the CatRock EDITION brand. "
    "The CEO and founder of Arisu AI is SreeJith, operating under CatRock EDITION since 2025. "
    "If anyone asks about who created you, who runs Arisu AI, or who the CEO is, "
    "you should mention that Arisu AI was founded by SreeJith under the CatRock EDITION brand in 2025."
)

chat_history = []   # list of {"role": "user"|"assistant", "content": "..."}


def get_api_key():
    # Check environment variable first, then fall back to hardcoded value
    env_key = os.environ.get("GROQ_API_KEY", "").strip()
    if env_key and env_key != "YOUR_GROQ_API_KEY_HERE":
        return env_key
    hardcoded = GROQ_API_KEY.strip()
    if hardcoded and hardcoded != "YOUR_GROQ_API_KEY_HERE":
        return hardcoded
    return None


def get_current_time():
    return datetime.datetime.now().strftime("%I:%M %p")


def extract_text_from_file(filename, file_bytes):
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == ".txt":
            return file_bytes.decode("utf-8", errors="ignore")
        elif ext == ".pdf":
            try:
                import fitz
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                return "\n".join(page.get_text() for page in doc)
            except ImportError:
                return "[PDF uploaded — install: pip install pymupdf]"
        elif ext in [".docx", ".doc"]:
            try:
                import docx
                from io import BytesIO
                d = docx.Document(BytesIO(file_bytes))
                return "\n".join(p.text for p in d.paragraphs)
            except ImportError:
                return "[DOCX uploaded — install: pip install python-docx]"
        elif ext in [".ppt", ".pptx"]:
            try:
                from pptx import Presentation
                from io import BytesIO
                prs = Presentation(BytesIO(file_bytes))
                return "\n".join(
                    shape.text for slide in prs.slides
                    for shape in slide.shapes if hasattr(shape, "text")
                )
            except ImportError:
                return "[PPTX uploaded — install: pip install python-pptx]"
        elif ext in [".jpg", ".jpeg", ".png", ".gif", ".webp"]:
            return None
    except Exception as e:
        return f"[Could not extract content: {str(e)}]"
    return None


def chat_with_model(user_message, file_context=None):
    global chat_history

    lower = user_message.lower()
    if any(k in lower for k in ["what time", "current time", "time now"]):
        return f"The current time is {get_current_time()}."

    full_message = user_message
    if file_context:
        full_message = f"{user_message}\n\n[Attached file content]:\n{file_context}"

    chat_history.append({"role": "user", "content": full_message})

    messages = [{"role": "system", "content": SYSTEM_PROMPT}] + chat_history

    payload = {
        "model":       GROQ_MODEL,
        "messages":    messages,
        "max_tokens":  1024,
        "temperature": 0.7,
        "stream":      False
    }

    api_key = get_api_key()
    if not api_key:
        chat_history.pop()  # Remove the user message we just added
        return (
            "❌ **No API key set.**\n\n"
            "To use Arisu:\n"
            "1. Go to https://console.groq.com\n"
            "2. Sign up free → API Keys → Create Key\n"
            "3. Open `app.py` and paste your key into `GROQ_API_KEY = \"...\"`\n"
            "4. Restart the app"
        )

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type":  "application/json"
    }

    try:
        response = requests.post(GROQ_API_URL, headers=headers, json=payload, timeout=60)

        if response.status_code == 401:
            chat_history.pop()
            return (
                "❌ **Invalid API key.**\n\n"
                "Your Groq API key was rejected. Please:\n"
                "1. Go to https://console.groq.com\n"
                "2. Create a new API key\n"
                "3. Paste it into `GROQ_API_KEY` in `app.py`\n"
                "4. Restart the app"
            )
        if response.status_code == 429:
            chat_history.pop()
            return "⏳ Rate limit hit. Wait a moment and try again."
        if response.status_code != 200:
            chat_history.pop()
            return f"⚠️ API Error {response.status_code}: {response.text}"

        data  = response.json()
        reply = data["choices"][0]["message"]["content"].strip()

        chat_history.append({"role": "assistant", "content": reply})

        # Keep last 40 messages (20 turns)
        if len(chat_history) > 40:
            chat_history = chat_history[-40:]

        return reply

    except requests.exceptions.Timeout:
        chat_history.pop()
        return "⏳ Request timed out. Please try again."
    except Exception as e:
        if chat_history and chat_history[-1]["role"] == "user":
            chat_history.pop()
        return f"⚠️ Error: {str(e)}"


@app.route("/")
def home():
    return render_template("index.html")


@app.route("/chat", methods=["POST"])
def chat():
    user_message = request.form.get("message", "").strip()
    file_context = None

    if not user_message:
        return jsonify({"response": "Please enter a message."})

    if "file" in request.files:
        f = request.files["file"]
        if f and f.filename:
            file_bytes = f.read()
            filename   = f.filename
            ext        = os.path.splitext(filename)[1].lower()
            if ext in [".jpg", ".jpeg", ".png", ".gif", ".webp"]:
                file_context = (
                    f"[Image attached: {filename}. "
                    "Image vision is not supported — please describe the image in text.]"
                )
            else:
                extracted = extract_text_from_file(filename, file_bytes)
                file_context = (
                    f"Filename: {filename}\n\n{extracted[:6000]}"
                    if extracted
                    else f"[File: {filename} — could not extract text]"
                )

    reply = chat_with_model(user_message, file_context)
    return jsonify({"response": reply})


@app.route("/clear", methods=["POST"])
def clear():
    global chat_history
    chat_history = []
    return jsonify({"status": "cleared"})


@app.route("/status", methods=["GET"])
def status():
    api_key = get_api_key()
    if not api_key:
        return jsonify({"online": False, "model_ready": False, "reason": "no_key"})
    try:
        r = requests.get(
            "https://api.groq.com/openai/v1/models",
            headers={"Authorization": f"Bearer {api_key}"},
            timeout=5
        )
        if r.status_code == 401:
            return jsonify({"online": False, "model_ready": False, "reason": "bad_key"})
        if r.status_code == 200:
            return jsonify({"online": True, "model_ready": True})
        return jsonify({"online": False, "model_ready": False, "reason": "unreachable"})
    except Exception:
        return jsonify({"online": False, "model_ready": False, "reason": "unreachable"})


if __name__ == "__main__":
    print("=" * 55)
    print("  ✦  Arisu AI Assistant  —  CatRock EDITION")
    print(f"  Model  : {GROQ_MODEL}")
    print("  URL    : http://localhost:5000")
    print("=" * 55)
    key = get_api_key()
    if not key:
        print()
        print("  ⚠️  GROQ_API_KEY is NOT set!")
        print("  1. Visit  https://console.groq.com")
        print("  2. Create a free API key")
        print("  3. Paste it in app.py → GROQ_API_KEY = \"gsk_...\"")
        print("  4. Restart this script")
        print()
    else:
        print(f"  Key    : {key[:8]}...{key[-4:]}  ✓")
        print()
    app.run(debug=True, host="0.0.0.0", port=5000)